#!/usr/bin/env python3
"""批量同步PPTX术语到最新版本"""
import os, sys
sys.path.insert(0, '/home/xspace/miniconda3/envs/lerobot/lib/python3.12/site-packages')
from pptx import Presentation

replacements = {
    "SmolVLA": "VTLA",
    "VLA 视觉语言动作模型": "多模态视觉语言动作混合模型",
    "触觉闭环反馈": "视触觉闭环反馈",
    "LeWorldModel": "因果世界模型",
    "精细感知全自主": "精细感知全自主安全",
    "≥99.2%": "大于99%",
    "全流程自主插拔": "全流程自主进行光模块插拔",
}

docs_dir = "/home/xspace/lerobot-smolvla-lew/docs"
pptx_files = [
    "L1-Z-MAX产品发布-v1.0.4.pptx",
    "Z-MAX产品培训-L2基线版.pptx",
    "Z-MAX基线版交付物培训.pptx",
    "Z-MAX基线版硬件配置培训.pptx",
    "Z700F-L2产品培训手册.pptx",
]

for fname in pptx_files:
    path = os.path.join(docs_dir, fname)
    if not os.path.exists(path):
        print(f"SKIP {fname} — not found")
        continue
    
    prs = Presentation(path)
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for old, new in replacements.items():
                            if old in run.text:
                                run.text = run.text.replace(old, new)
                                count += 1
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                for old, new in replacements.items():
                                    if old in run.text:
                                        run.text = run.text.replace(old, new)
                                        count += 1
    
    prs.save(path)
    print(f"✅ {fname}: {count} replacements")

print("\nAll PPTX updated!")
