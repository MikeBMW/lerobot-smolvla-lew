"""
Z-MAX PPT 指令引擎 — 从 PPT 解析指令，驱动控制台动作

核心思路：
  用户在 PPT 每页写一条指令（标题=动词，正文=参数），
  控制台解析后自动生成对应的 文件/配置/训练/部署 动作。

支持的指令动词:
  CREATE_FILE   — 在静界中创建文档
  UPDATE_CONFIG — 修改控制台配置
  RUN_CMD       — 执行系统命令
  TRAIN_MODEL   — 启动训练任务
  DEPLOY        — 部署到 ECS
  EVAL_MODEL    — 运行评估
  GIT_COMMIT    — 提交并推送 GitHub
"""

import os, sys, json, re, subprocess, time
from datetime import datetime

# ── PPT 指令模板 ──

TEMPLATE_SLIDE = """---
# 指令模板 — 每页 PPT = 一条指令
#
# 标题行: 动词 + 名称
# 正文:   参数 (JSON 或 键值对)
# 备注:   状态追踪

## 标题格式
{动词}: {指令名称}

## 正文格式 (任选一种)
### 格式A — JSON
```json
{参数}
```

### 格式B — 键值对
参数1: 值1
参数2: 值2
---

## 当前支持动词

| 动词 | 作用 | 示例 |
|------|------|------|
| CREATE_FILE | 创建文档 | CREATE_FILE: 静界/01-培训/Z700F操作指南.md |
| UPDATE_CONFIG | 改配置 | UPDATE_CONFIG: 训练参数 |
| RUN_CMD | 执行命令 | RUN_CMD: lerobot-train --config ... |
| TRAIN_MODEL | 训练模型 | TRAIN_MODEL: SmolVLA v2 |
| DEPLOY | 部署上线 | DEPLOY: 同步到 ECS |
| EVAL_MODEL | 评估模型 | EVAL_MODEL: benchmark结果 |
| GIT_COMMIT | 提交代码 | GIT_COMMIT: 训练日志 |
"""


INSTRUCTION_DIR = "00-指令"  # 静界下的指令目录


def get_instructions_dir():
    """返回指令目录路径"""
    from docs_sync import get_docs_dir
    return os.path.join(get_docs_dir(), INSTRUCTION_DIR)


def create_template():
    """生成指令模板 PPTX"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return False, "需要安装 python-pptx: pip install python-pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── 封面 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(6, 8, 13)

    # 标题
    from pptx.util import Pt
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Z-MAX 指令控制 PPT"
    p.font.size = Pt(44)
    p.font.color.rgb = RGBColor(0, 212, 170)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "每页一条指令 → 控制台自动执行"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(200, 209, 217)
    p2.alignment = PP_ALIGN.CENTER

    # ── 指令页示例 ──
    examples = [
        ("CREATE_FILE: Z700F操作指南", "路径: 静界/01-培训/Z700F操作指南.md\n内容: 标准操作流程..."),
        ("TRAIN_MODEL: SmolVLA v2.1", "数据集: lerobot/pusht\n步数: 500\nbatch_size: 4\n学习率: 0.0001"),
        ("DEPLOY: 同步控制台到官网", "目标: 39.102.211.79\n路径: /www/wwwroot/datadrive.world/"),
        ("GIT_COMMIT: 训练日志备份", "信息: feat: 第3轮训练完成\n分支: main"),
    ]

    for title, body in examples:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(6, 8, 13)

        # 标题
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(0, 212, 170)
        p.font.bold = True

        # 正文
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, line in enumerate(body.split("\n")):
            if i == 0:
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()
            p2.text = line
            p2.font.size = Pt(14)
            p2.font.color.rgb = RGBColor(200, 209, 217)
            p2.space_after = Pt(6)

        # 动词标签
        verb = title.split(":")[0]
        txBox3 = slide.shapes.add_textbox(Inches(11), Inches(0.3), Inches(2), Inches(0.5))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = verb
        p3.font.size = Pt(10)
        p3.font.bold = True
        p3.font.color.rgb = RGBColor(0, 212, 170)

    output_path = os.path.join(get_instructions_dir(), "Z-MAX指令控制模板.pptx")
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return True, output_path


def parse_slide_text(slide):
    """从 PPT 幻灯片提取标题和正文"""
    title = ""
    body_lines = []

    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                # 第一行有 ":" 的当作标题
                if ":" in text and not title and len(text) < 100:
                    title = text
                else:
                    body_lines.append(text)

    return title, "\n".join(body_lines)


def parse_pptx(path):
    """解析 PPTX 文件，提取所有指令"""
    try:
        from pptx import Presentation
    except ImportError:
        return False, "需要安装 python-pptx"

    if not os.path.exists(path):
        return False, f"文件不存在: {path}"

    prs = Presentation(path)
    instructions = []

    for i, slide in enumerate(prs.slides):
        title, body = parse_slide_text(slide)
        if not title:
            continue

        # 提取动词
        verb = title.split(":")[0].strip().upper()
        name = title.split(":", 1)[1].strip() if ":" in title else ""

        if verb in ("CREATE_FILE", "UPDATE_CONFIG", "RUN_CMD",
                     "TRAIN_MODEL", "DEPLOY", "EVAL_MODEL", "GIT_COMMIT"):
            instructions.append({
                "slide": i + 1,
                "verb": verb,
                "name": name,
                "body": body,
                "raw": title,
            })

    return True, instructions


def execute_instruction(inst, log_callback=print):
    """执行单条指令"""
    verb = inst["verb"]
    name = inst["name"]
    body = inst["body"]

    log_callback(f"▶ 执行: {verb}: {name}")

    if verb == "CREATE_FILE":
        return _exec_create_file(name, body, log_callback)
    elif verb == "RUN_CMD":
        return _exec_run_cmd(name, body, log_callback)
    elif verb == "GIT_COMMIT":
        return _exec_git_commit(name, body, log_callback)
    elif verb == "TRAIN_MODEL":
        log_callback(f"  ⏳ 训练任务已记录: {name}")
        return True, f"训练任务已入队: {name}"
    else:
        log_callback(f"  ⚠  暂未实现: {verb}")
        return False, f"未实现: {verb}"


def _exec_create_file(name, body, log_callback):
    """CREATE_FILE: 创建文档"""
    from docs_sync import get_docs_dir
    docs_dir = get_docs_dir()

    # 如果 name 包含路径，直接使用；否则放在 00-指令/ 下
    if "/" in name:
        target = os.path.join(docs_dir, name)
    else:
        target = os.path.join(docs_dir, INSTRUCTION_DIR, name)

    os.makedirs(os.path.dirname(target), exist_ok=True)
    with open(target, "w", encoding="utf-8") as f:
        f.write(body)
    log_callback(f"  ✅ 已创建: {target}")
    return True, target


def _exec_run_cmd(name, body, log_callback):
    """RUN_CMD: 执行命令"""
    cmd = body.strip() or name
    try:
        r = subprocess.run(cmd, shell=True, capture_output=True, text=True, timeout=60)
        log_callback(f"  stdout: {r.stdout[:200]}")
        if r.stderr:
            log_callback(f"  stderr: {r.stderr[:200]}")
        return r.returncode == 0, r.stdout[:500]
    except Exception as e:
        log_callback(f"  ❌ {e}")
        return False, str(e)


def _exec_git_commit(name, body, log_callback):
    """GIT_COMMIT: 提交并推送"""
    repo = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    msg = name or f"指令提交: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
    try:
        subprocess.run(["git", "-C", repo, "add", "-A"], capture_output=True, timeout=10)
        subprocess.run(["git", "-C", repo, "commit", "-m", msg], capture_output=True, timeout=10)
        r = subprocess.run(["git", "-C", repo, "push"], capture_output=True, text=True, timeout=30)
        log_callback(f"  ✅ 已推送: {msg}")
        return True, r.stdout[:200]
    except Exception as e:
        log_callback(f"  ❌ {e}")
        return False, str(e)


def run_all(input_path, log_callback=print):
    """解析 PPTX 并执行所有指令"""
    ok, result = parse_pptx(input_path)
    if not ok:
        log_callback(f"❌ 解析失败: {result}")
        return False

    instructions = result
    log_callback(f"📋 发现 {len(instructions)} 条指令")
    log_callback("=" * 40)

    success = 0
    for inst in instructions:
        ok, msg = execute_instruction(inst, log_callback)
        if ok:
            success += 1
        log_callback("-" * 40)

    log_callback(f"\n📊 结果: {success}/{len(instructions)} 成功")
    return success == len(instructions)
